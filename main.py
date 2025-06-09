import os
from flask import Flask, request, render_template_string, jsonify, send_file, url_for
import PyPDF2
import re
import requests
import json
from io import BytesIO
from pptx import Presentation # Library for .pptx files
from pptx.enum.shapes import MSO_SHAPE_TYPE # For checking shape type
from docx import Document
from docx.shared import Inches, Pt # Added Pt for font size
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.dml import MSO_THEME_COLOR
from docx.opc.constants import RELATIONSHIP_TYPE as RT # For embedding images in DOCX
from PIL import Image # For converting image formats if needed, and saving

# NEW IMPORTS FOR SEMANTIC SEARCH
from sentence_transformers import SentenceTransformer, util
import torch # SentenceTransformer uses PyTorch or TensorFlow, ensuring torch is available if not already.


app = Flask(__name__)

# --- Configuration ---
# IMPORTANT: Set your Gemini API key here or, even better, as an environment variable.
# For example: export GEMINI_API_KEY="YOUR_API_KEY_HERE"
# Using os.environ.get to allow setting via environment variable, with a placeholder fallback.
GEMINI_API_KEY = ''
GEMINI_API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash:generateContent?key={GEMINI_API_KEY}"

# Firebase configuration variables (kept as they were in your previous code, but not actively used
# for data persistence in this version as per your request).
FIREBASE_CONFIG_JSON_STR = os.environ.get('__firebase_config', '{}')
try:
    firebase_config = json.loads(FIREBASE_CONFIG_JSON_STR)
except json.JSONDecodeError:
    firebase_config = {} # Default to empty if parsing fails

# In-memory storage for uploaded content.
# This data will be lost when the server restarts.
QUESTION_BANK_CONTENT = ""
QUESTION_BANK_FILENAME = ""
SLIDE_CONTENTS = {} # Stores {filename: full_text_content}

# Directory for temporarily storing extracted images
TEMP_IMAGE_DIR = os.path.join(app.root_path, 'static', 'temp_images')


# Initialize SentenceTransformer model for semantic search globally
# This ensures the model is loaded only once when the Flask app starts.
print("Loading SentenceTransformer model...")
try:
    # 'all-MiniLM-L6-v2' is a good balance of speed and performance.
    # For higher accuracy, consider 'all-mpnet-base-v2' (larger model).
    sentence_transformer_model = SentenceTransformer('all-MiniLM-L6-v2')
    print("SentenceTransformer model loaded successfully.")
except Exception as e:
    print(f"Error loading SentenceTransformer model: {e}")
    print("Semantic search will be disabled. Ensure you have internet access on first run to download the model.")
    sentence_transformer_model = None # Set to None if loading fails

# Ensure the temporary image directory exists
if not os.path.exists(TEMP_IMAGE_DIR):
    os.makedirs(TEMP_IMAGE_DIR)

# --- Helper Functions ---

def extract_text_from_pdf(pdf_file_stream):
    """
    Extracts text from a PDF file stream.
    
    Args:
        pdf_file_stream: A file-like object (BytesIO) containing the PDF data.
        
    Returns:
        A string containing the extracted text from all pages.
    """
    text = ""
    try:
        reader = PyPDF2.PdfReader(pdf_file_stream)
        for page in reader.pages:
            text += page.extract_text() or "" # Ensure non-None text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_pptx(pptx_file_stream):
    """
    Extracts text from a PPTX file stream. (Legacy for slides without image embedding)
    
    Args:
        pptx_file_stream: A file-like object (BytesIO) containing the PPTX data.
        
    Returns:
        A string containing the concatenated text from all slides.
    """
    text = ""
    try:
        prs = Presentation(pptx_file_stream)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text += shape.text_frame.text + "\n"
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
    return text

def extract_content_from_pptx(pptx_file_stream, temp_image_dir=TEMP_IMAGE_DIR):
    """
    Extracts text and embeds HTML image tags for pictures from a PPTX file.
    Saves images to a temporary directory and replaces [Image X] placeholders.
    """
    full_text_with_images = ""
    
    if not os.path.exists(temp_image_dir):
        os.makedirs(temp_image_dir)

    try:
        prs = Presentation(pptx_file_stream)
        for slide_idx, slide in enumerate(prs.slides):
            current_slide_text_parts = []
            current_slide_images = [] # Store extracted image info for this slide

            # First pass: Extract text and images
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    current_slide_text_parts.append(shape.text_frame.text)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_bytes = shape.image.blob
                    image_extension = shape.image.ext # 'png', 'jpeg' etc.
                    
                    # Ensure image_extension is valid for PIL
                    if image_extension == 'wmf': # PIL doesn't support WMF directly without external libs
                        try:
                            # Attempt conversion for unsupported formats like WMF
                            img = Image.open(BytesIO(image_bytes))
                            image_bytes = BytesIO()
                            img.save(image_bytes, format='PNG')
                            image_extension = 'png'
                            image_bytes.seek(0)
                        except Exception as convert_e:
                            print(f"Warning: Could not convert WMF image. Skipping. {convert_e}")
                            continue # Skip this image if conversion fails

                    # Create a unique filename for the image
                    # Use a hash of image bytes to avoid saving duplicates for identical images
                    image_hash = hashlib.sha256(image_bytes).hexdigest() if isinstance(image_bytes, bytes) else hashlib.sha256(image_bytes.getvalue()).hexdigest()
                    image_name = f"slide_{slide_idx}_{image_hash[:8]}.{image_extension}"
                    image_full_path = os.path.join(temp_image_dir, image_name)
                    image_url = f"/static/temp_images/{image_name}" # URL for frontend

                    # Save the image to the temporary directory
                    if not os.path.exists(image_full_path): # Only save if not already saved
                        if isinstance(image_bytes, BytesIO):
                            with open(image_full_path, "wb") as f:
                                f.write(image_bytes.getvalue())
                        else: # Assume it's already bytes
                             with open(image_full_path, "wb") as f:
                                f.write(image_bytes)

                    current_slide_images.append({
                        'id': image_name,
                        'url': image_url
                    })
            
            combined_text = "\n".join(current_slide_text_parts)

            # Second pass: Replace [Image X] placeholders with HTML <img> tags
            # The placeholders are typically found in the text itself.
            # Assuming [Image 1], [Image 2], etc. corresponds to the order of extraction.
            # This is a heuristic and might not be perfect for all PPTXs.
            for img_idx, img_info in enumerate(current_slide_images):
                # We need to be careful with the exact placeholder text.
                # From 1. Auto Encoders.pptx, it looks like "[Image 14]"
                placeholder_regex = re.compile(r'\[Image\s*' + str(img_idx + 1) + r'\]')
                # Replace the placeholder with the HTML <img> tag
                combined_text = placeholder_regex.sub(
                    f"<img src='{img_info['url']}' alt='Slide Image {img_idx + 1}' class='extracted-image' style='max-width: 100%; height: auto; border-radius: 8px; margin: 10px 0;'>",
                    combined_text,
                    count=1 # Replace only the first occurrence to avoid issues if multiple placeholders exist but only one image is intended
                )
            
            full_text_with_images += combined_text.strip() + "\n\n-- PAGE " + str(slide_idx + 1) + " --\n\n" # Add page break for clarity
            
    except Exception as e:
        print(f"Error extracting content from PPTX: {e}")
        # Optionally, clean up partially created temp_images here if an error occurs
    
    return full_text_with_images.strip() # Return the combined text with embedded image tags

def find_relevant_snippets_semantic(query_text, all_slide_texts, top_n_snippets=3):
    """
    Finds semantically relevant snippets from lecture slides based on a query using sentence embeddings.
    If the SentenceTransformer model is not loaded, it falls back to keyword search.
    
    Args:
        query_text (str): The question text to find relevant content for.
        all_slide_texts (dict): A dictionary where keys are filenames and values are the
                                full text content of each slide file.
        top_n_snippets (int): Maximum number of top relevant snippets to return.

    Returns:
        list: A list of dictionaries, each containing 'filename' and 'snippet'.
    """
    # If the semantic model failed to load, use keyword fallback directly
    if sentence_transformer_model is None:
        print("SentenceTransformer model not loaded. Falling back to simple keyword search.")
        return find_relevant_snippets_keyword(query_text, all_slide_texts, top_n_snippets)

    query_embedding = sentence_transformer_model.encode(query_text, convert_to_tensor=True)
    all_scored_snippets = []

    for filename, text_content in all_slide_texts.items():
        # Temporarily remove <img> tags for embedding consistency, they will be back in the answer
        cleaned_text_content = re.sub(r'<img[^>]*>', '', text_content)
        
        # Split text into paragraphs/chunks for better granularity
        paragraphs = re.split(r'\n{2,}', cleaned_text_content)
        # Filter out very short or empty paragraphs to avoid noisy embeddings
        relevant_paragraphs = [p.strip() for p in paragraphs if p.strip() and len(p.strip()) > 20]

        if not relevant_paragraphs:
            continue

        # Generate embeddings for all relevant paragraphs
        paragraph_embeddings = sentence_transformer_model.encode(relevant_paragraphs, convert_to_tensor=True)

        # Calculate cosine similarity between query embedding and all paragraph embeddings
        cosine_scores = util.pytorch_cos_sim(query_embedding, paragraph_embeddings)[0]

        for i, score in enumerate(cosine_scores):
            all_scored_snippets.append({
                'score': score.item(), # Convert PyTorch tensor item to a Python number
                'snippet': relevant_paragraphs[i],
                'filename': filename
            })

    # Sort snippets by similarity score in descending order
    all_scored_snippets.sort(key=lambda x: x['score'], reverse=True)

    final_snippets_data = []
    unique_snippet_texts = set() # To prevent adding duplicate snippets
    for item in all_scored_snippets:
        if item['snippet'] not in unique_snippet_texts:
            final_snippets_data.append({'filename': item['filename'], 'snippet': item['snippet']})
            unique_snippet_texts.add(item['snippet'])
        if len(final_snippets_data) >= top_n_snippets:
            break # Stop once we have enough top snippets

    if not final_snippets_data:
        print("Semantic search yielded no results. Attempting keyword fallback.")
        return find_relevant_snippets_keyword(query_text, all_slide_texts, top_n_snippets)

    return final_snippets_data


def find_relevant_snippets_keyword(query, all_slide_texts, top_n_snippets=3):
    """
    Finds relevant content snippets from slide texts based on keyword matching (fallback).
    """
    query_words = set(re.findall(r'\b\w+\b', query.lower()))
    all_scored_snippets = []

    for filename, text_content in all_slide_texts.items():
        # Temporarily remove <img> tags for cleaner keyword search
        cleaned_text_content = re.sub(r'<img[^>]*>', '', text_content)

        paragraphs = re.split(r'\n{2,}', cleaned_text_content) 
        
        for para in paragraphs:
            para_lower = para.lower()
            score = len(query_words.intersection(set(re.findall(r'\b\w+\b', para_lower))))
            
            if score > 0 and para.strip():
                all_scored_snippets.append({'score': score, 'snippet': para.strip(), 'filename': filename})

    all_scored_snippets.sort(key=lambda x: x['score'], reverse=True)
    
    final_snippets_data = []
    unique_snippet_texts = set()
    for item in all_scored_snippets:
        if item['snippet'] not in unique_snippet_texts:
            final_snippets_data.append({'filename': item['filename'], 'snippet': item['snippet']})
            unique_snippet_texts.add(item['snippet'])
        if len(final_snippets_data) >= top_n_snippets:
            break
            
    return final_snippets_data


def parse_questions_from_text(full_text):
    """
    Parses questions from the full extracted text, specifically handling
    'QXXX (a)', 'QXXX (b)', and ' (Or) ' clauses to extract distinct questions.
    Assigns sequential Q IDs. Preserves <img> tags.
    """
    questions = []
    current_q_id = 1
    
    # Pre-clean the text to remove common headers/footers/etc.
    # IMPORTANT: Ensure HTML <img> tags are NOT stripped by these regexes.
    # The pattern will explicitly match the image tag and NOT replace it.
    
    # Temporarily hide img tags from general cleaning
    img_placeholder_map = {}
    def replace_img_with_placeholder(match):
        img_tag = match.group(0)
        placeholder = f"__IMG_PLACEHOLDER_{len(img_placeholder_map)}__"
        img_placeholder_map[placeholder] = img_tag
        return placeholder

    text_with_placeholders = re.sub(r'<img[^>]*>', replace_img_with_placeholder, full_text)


    excluded_line_patterns = [
        r'^\s*Common to \(AI-DS, AI-ML\).*$',
        r'^\s*19AI413– DEEP LEARNING AND ITS APPLICATIONS.*$',
        r'^\s*Time:\s*Three\s+hours.*$',
        r'^\s*Maximum\s+marks:\s*\d+.*$',
        r'^\s*Question Repository.*$',
        r'^\s*Faculty Name.*$', r'^\s*Department AIDS.*$', # Ensure these don't accidentally contain image info
        r'^\s*Answer All Questions.*$',
        r'^\s*PART\s+[A-C]\s*\(.*\).*$', # Matches "PART A (10 x 2 = 20 marks)"
        r'^\s*Q\.\s*No\s*Questions\s*CO\s*Knowledge\s*Level\s*\(Blooms\)\s*Difficulty\s*Level\s*\(1-5\)\s*$', # Table header
        r'^\s*CO\s*Knowledge\s*Level\s*Difficulty\s*Level.*$', # Another table header
        r'^\s*\(Blooms\)\s*\(1-5\)\s*$', # Just Blooms and Difficulty levels
        r'^\s*\d+\s*$', # Isolated numbers (e.g., page numbers, or just "3" etc.)
        r'^\s*x\d+\s*$', r'^\s*w\d+\s*$', r'^\s*b\s*$', r'^\s*\+\s*$', r'^\s*y\s*$', # Single char/num lines from diagrams
        r'^\s*Input\s*$', r'^\s*output\s*$', r'^\s*hidden\s+layer\s+\d+\s*$', r'^\s*softmax\s*$', r'^\s*Linear\s*$',
        r'^\s*CO\d\s*$', r'^\s*K\d\s*$', r'^\s*Knowledge\s+Level\s*$', r'^\s*Difficulty\s+Level\s*$',
        r'^\s*Reg\. No\s*$', r'^\s*QP\s+Code\s*$',
        r'^\s*(?:Note:)\s*.*$', # Lines starting with "Note:"
        r'^\s*-- PAGE\s+\d+\s*--$', # Page break lines
        r'^\s*\(Case study/Comprehensive type Questions\)\s*$',
        r'^\s*\d+\s*x\s*\d+\s*=\s*\d+\s*marks\)\s*$', # "(X x Y = Z marks)"
        r'^\s*Questions\s*$', # Just "Questions"
        r'^\s*\(i+\)\s*$', r'^\s*\(a\)\s*$', r'^\s*\(b\)\s*$', # Isolated sub-question markers
        r'^\s*\(Or\)\s*$', # Isolated (Or)
        r'^\s*-\s*$' # Isolated hyphen
    ]

    lines = text_with_placeholders.split('\n')
    cleaned_lines = []
    for line in lines:
        stripped_line = line.strip()
        if not stripped_line: # Remove empty lines
            continue
        
        is_excluded = False
        for pattern in excluded_line_patterns:
            if re.fullmatch(pattern, stripped_line, re.IGNORECASE):
                is_excluded = True
                break
        
        if not is_excluded:
            # Further clean up any inline CO/K codes or Marks
            stripped_line = re.sub(r'\s*\bCO\d\b\s*|\s*\bK\d\b\s*', '', stripped_line, flags=re.IGNORECASE).strip()
            stripped_line = re.sub(r'\s*\(\d+\s*Marks?\)\s*', '', stripped_line, flags=re.IGNORECASE).strip()
            stripped_line = re.sub(r'\s*\d+\s*$', '', stripped_line).strip() # Trailing numbers like "3" or "4"
            
            if stripped_line: # Add if not empty after all cleanups
                cleaned_lines.append(stripped_line)

    # Rejoin cleaned lines with double newlines to form blocks, then split into major question blocks
    cleaned_text_with_placeholders = '\n\n'.join(cleaned_lines)
    
    # Pattern to identify the start of a new major question (QA, QB, QC)
    major_q_start_pattern = re.compile(r'^(Q[ABC]\d{3}(?:\s*\(?[ab]\)?)?)', re.MULTILINE)
    
    split_parts = major_q_start_pattern.split(cleaned_text_with_placeholders)

    content_blocks = [part.strip() for part in split_parts if part.strip()]

    full_q_blocks = []
    i = 0
    while i < len(content_blocks):
        if major_q_start_pattern.match(content_blocks[i]):
            current_id = content_blocks[i]
            current_content = ""
            i += 1
            while i < len(content_blocks) and not major_q_start_pattern.match(content_blocks[i]):
                current_content += "\n\n" + content_blocks[i]
                i += 1
            full_q_blocks.append((current_id, current_content.strip()))
        else:
            i += 1

    for original_q_id_prefix, block_content_with_placeholders in full_q_blocks:
        or_parts = re.split(r'\s*\b(?:Or|or)\b\s*', block_content_with_placeholders, flags=re.IGNORECASE)
        
        for or_part_index, or_sub_block_with_placeholders in enumerate(or_parts):
            sub_q_pattern = re.compile(r'(\([a-zivx]+\))', re.IGNORECASE)
            sub_q_parts = sub_q_pattern.split(or_sub_block_with_placeholders)
            
            temp_q_text_with_placeholders = ""
            current_sub_id_suffix = ""
            
            if re.match(r'^\([a-zivx]+\)$', original_q_id_prefix):
                 current_sub_id_suffix = original_q_id_prefix

            for part_idx, part in enumerate(sub_q_parts):
                part_stripped = part.strip()
                if not part_stripped:
                    continue
                
                is_sub_q_marker = re.match(r'^\([a-zivx]+\)$', part_stripped, re.IGNORECASE)
                
                if is_sub_q_marker:
                    if temp_q_text_with_placeholders:
                        final_question_text = re.sub(r'\s*\d+$', '', temp_q_text_with_placeholders).strip()
                        if final_question_text:
                            # Restore image tags before adding to questions list
                            for ph, img_tag in img_placeholder_map.items():
                                final_question_text = final_question_text.replace(ph, img_tag)
                            questions.append({
                                'id': f'Q{current_q_id}{current_sub_id_suffix}',
                                'text': final_question_text
                            })
                            current_q_id += 1
                    current_sub_id_suffix = part_stripped
                    temp_q_text_with_placeholders = ""
                else:
                    temp_q_text_with_placeholders += " " + part_stripped
            
            if temp_q_text_with_placeholders.strip():
                final_question_text = re.sub(r'\s*\d+$', '', temp_q_text_with_placeholders).strip()
                if final_question_text:
                    # Restore image tags before adding to questions list
                    for ph, img_tag in img_placeholder_map.items():
                        final_question_text = final_question_text.replace(ph, img_tag)
                    questions.append({
                        'id': f'Q{current_q_id}{current_sub_id_suffix}',
                        'text': final_question_text
                    })
                    current_q_id += 1
                temp_q_text_with_placeholders = ""

    for q in questions:
        # Final cleanup pass, specifically for remnants of Q IDs and numbers
        q['text'] = re.sub(r'^(Q[ABC]\d{3}(?:\s*\(?[ab]\)?)?\s*|\([a-zivx]+\)\s*)', '', q['text'], flags=re.IGNORECASE).strip()
        q['text'] = re.sub(r'^\s*\d+\.?\s*', '', q['text']).strip() # Remove leading numbers like "1. "
        
        # Restore any image tags that might have been hidden/modified
        for ph, img_tag in img_placeholder_map.items():
            q['text'] = q['text'].replace(ph, img_tag)

    return questions


def prepare_gemini_prompt(question_text, context_text):
    """
    Prepares the prompt for the Gemini API, including the question and relevant context.
    The prompt is designed to instruct Gemini on content ratio and highlighting.
    """
    # Remove image tags from the question text sent to Gemini, as the LLM doesn't process images visually
    clean_question_text = re.sub(r'<img[^>]*>', '', question_text)

    prompt = f"""You are an expert academic assistant specializing in Deep Learning.
    
    Your task is to answer the following question.
    
    **Question:** {clean_question_text}
    
    **Instructions:**
    1.  **Prioritize content from the provided lecture slides (approximately 75% of the answer).**
    2.  **Supplement with your general knowledge (approximately 25% of the answer).**
    3.  **Highlight key points using bold text (e.g., **important concept**) or bullet points (e.g., * list item).**
    4.  **Structure your answer clearly and concisely.**
    5.  **If information from slides is available, ensure it forms the core of your response.**
    
    **Relevant Information from Lecture Slides:**
    ---
    {context_text}
    ---
    
    Please provide a comprehensive answer to the question, adhering to the above instructions.
    """
    return prompt


# --- Flask Routes ---

@app.route('/')
def index():
    """
    Renders the main HTML page for the application.
    Passes environment-injected variables to the frontend.
    """
    # These variables are expected to be set by the Canvas environment for Firebase setup.
    # If Firebase is not being used for persistence, these can be default empty strings.
    app_id = os.environ.get('__app_id', 'default-app-id')
    initial_auth_token = os.environ.get('__initial_auth_token', '')

    # The HTML template is now part of this app.py file
    return render_template_string(HTML_TEMPLATE,
                                  app_id=json.dumps(app_id),
                                  firebase_config_json=json.dumps(firebase_config), # Pass the loaded config
                                  initial_auth_token=json.dumps(initial_auth_token))


@app.route('/upload_qb', methods=['POST'])
def upload_qb():
    """
    Handles the upload of a question bank file (PDF or TXT).
    Extracts text and parses questions, storing them in-memory.
    """
    global QUESTION_BANK_CONTENT, QUESTION_BANK_FILENAME
    if 'qb_file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    file = request.files['qb_file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    filename = file.filename
    file_extension = filename.split('.')[-1].lower()
    file_stream = BytesIO(file.read())
    
    text_content = ""
    if file_extension == 'pdf':
        text_content = extract_text_from_pdf(file_stream)
    elif file_extension == 'pptx': # Added PPTX support for QB with image extraction
        text_content = extract_content_from_pptx(file_stream) # Use new function for PPTX
    elif file_extension == 'txt':
        text_content = file_stream.getvalue().decode('utf-8')
    else:
        return jsonify({'error': 'Unsupported file type. Only PDF, PPTX, and TXT are supported for Question Banks.'}), 400
    
    if not text_content:
        return jsonify({'error': 'Could not extract text from the uploaded question bank.'}), 400

    QUESTION_BANK_CONTENT = text_content
    QUESTION_BANK_FILENAME = filename
    
    questions = parse_questions_from_text(text_content)
    
    return jsonify({
        'message': 'Question bank uploaded and parsed successfully!',
        'filename': filename,
        'questions': questions
    })

@app.route('/upload_slides', methods=['POST'])
def upload_slides():
    """
    Handles the upload of lecture slide files (PDF, PPTX, or TXT).
    Extracts text from each file and stores it in-memory.
    """
    global SLIDE_CONTENTS
    SLIDE_CONTENTS = {} # Clear previous slides when new ones are uploaded
    
    if 'slide_files' not in request.files:
        return jsonify({'error': 'No files part'}), 400
    
    files = request.files.getlist('slide_files')
    if not files:
        return jsonify({'error': 'No selected files'}), 400

    uploaded_count = 0
    processed_files_info = [] # To send back details about each processed file
    
    for file in files:
        if file.filename == '':
            continue
        
        filename = file.filename
        file_extension = filename.split('.')[-1].lower()
        file_stream = BytesIO(file.read())
        
        text_content = ""
        try:
            if file_extension == 'pdf':
                text_content = extract_text_from_pdf(file_stream)
            elif file_extension == 'pptx':
                # For slides, we still extract text for semantic search
                # But actual image embedding is more relevant if QB is PPTX
                # Here, we just extract text, images are saved but not embedded into the context passed to LLM yet.
                # The find_relevant_snippets_semantic/keyword handles image tags by stripping them for LLM.
                text_content = extract_content_from_pptx(file_stream) # Use new function for PPTX slides as well
            elif file_extension == 'txt':
                text_content = file_stream.getvalue().decode('utf-8')
            else:
                processed_files_info.append({'filename': filename, 'status': 'unsupported', 'error': 'Unsupported file type.'})
                continue # Skip to next file if type is not supported

            if not text_content:
                processed_files_info.append({'filename': filename, 'status': 'failed', 'error': 'Could not extract text.'})
                continue # Skip if no text extracted
            
            SLIDE_CONTENTS[filename] = text_content
            uploaded_count += 1
            processed_files_info.append({'filename': filename, 'status': 'success', 'content_length': len(text_content)})

        except Exception as e:
            print(f"Error processing slide file {filename}: {e}")
            processed_files_info.append({'filename': filename, 'status': 'error', 'error': str(e)})


    if not SLIDE_CONTENTS:
        return jsonify({'error': 'No valid lecture slide files were uploaded or text extracted. Supported types: PDF, PPTX, TXT.', 'processed_files': processed_files_info}), 400

    return jsonify({
        'message': f'Successfully processed {uploaded_count} lecture slide(s).',
        'processed_files': processed_files_info
    })


@app.route('/generate_answer', methods=['POST'])
def generate_answer():
    """
    Generates an answer for a single question using the Gemini API.
    It retrieves relevant context from the in-memory lecture slides using semantic search.
    """
    data = request.json
    question_text = data.get('question') # This text might contain <img> tags
    
    # --- Input Validation ---
    if not question_text:
        return jsonify({'error': 'Question text is required'}), 400
    if not SLIDE_CONTENTS:
        return jsonify({'error': 'No lecture slides uploaded. Please upload slides to generate answers.'}), 400
    if not GEMINI_API_KEY or GEMINI_API_KEY == 'YOUR_GEMINI_API_KEY_HERE':
        return jsonify({'error': 'Gemini API key is not configured on the server. Please set it.'}), 500

    # --- Context Retrieval (Semantic Search) ---
    # Use the globally stored slide contents for context
    # Note: relevant_snippets functions will strip <img> tags before sending to LLM.
    relevant_snippets = find_relevant_snippets_semantic(question_text, SLIDE_CONTENTS, top_n_snippets=3)
    
    context_text = "No directly relevant information found in lecture slides. Answering primarily based on general knowledge."
    source_references = [] # To store filename references

    if relevant_snippets:
        formatted_snippets = []
        for snippet_info in relevant_snippets:
            formatted_snippets.append(f"**From {snippet_info['filename']}**:\n{snippet_info['snippet']}")
            source_references.append(snippet_info['filename'])
        context_text = "## Relevant Information from Lecture Slides:\n" + "\n\n".join(formatted_snippets)
    
    # Ensure sources are unique and ordered alphabetically for consistent display
    source_references = sorted(list(set(source_references)))
    
    # --- Prepare LLM Prompt ---
    # The prepare_gemini_prompt function will strip <img> tags from question_text before sending to LLM.
    prompt_content = prepare_gemini_prompt(question_text, context_text)

    # --- Call Gemini API ---
    try:
        headers = {
            'Content-Type': 'application/json'
        }
        payload = {
            'contents': [
                {'parts': [{'text': prompt_content}]}
            ],
            "generationConfig": {
                "temperature": 0.2, # Lower temperature for more factual and less creative answers
                "topP": 0.9,
                "topK": 40,
                "maxOutputTokens": 1500 # Adjust for desired max answer length
            }
        }
        
        # Make the API request with a timeout
        response = requests.post(GEMINI_API_URL, headers=headers, json=payload, timeout=90) # Increased timeout
        response.raise_for_status() # Raise an exception for HTTP errors (4xx or 5xx)
        
        gemini_response = response.json()
        
        # Check if 'candidates' and 'content' exist before accessing
        if gemini_response and 'candidates' in gemini_response and gemini_response['candidates']:
            first_candidate = gemini_response['candidates'][0]
            if 'content' in first_candidate and 'parts' in first_candidate['content'] and first_candidate['content']['parts']:
                answer_text = first_candidate['content']['parts'][0].get('text', 'No answer generated.')
                
                # Append sources to the answer text
                sources_str = ", ".join(source_references) if source_references else "General knowledge (No specific slide content found)."
                full_answer = f"{answer_text}\n\n**Sources:** {sources_str}"

                return jsonify({'answer': full_answer})
            else:
                return jsonify({'error': 'Invalid response format from Gemini API (missing content parts).', 'details': gemini_response}), 500
        else:
            return jsonify({'error': 'Invalid response format from Gemini API (no candidates).', 'details': gemini_response}), 500

    except requests.exceptions.RequestException as e:
        # Handle network-related errors (e.g., connection refused, timeout)
        print(f"Error calling Gemini API: {e}")
        return jsonify({'error': f'Failed to get response from Gemini API: {e}'}), 500
    except json.JSONDecodeError as e:
        # Handle cases where the API response is not valid JSON
        print(f"Error decoding JSON response from Gemini API: {e}")
        return jsonify({'error': f'Invalid JSON response from Gemini API: {e}'}), 500
    except Exception as e:
        # Catch any other unexpected errors
        print(f"An unexpected error occurred: {e}")
        return jsonify({'error': f'An unexpected error occurred: {e}'}), 500


@app.route('/download_answers_docx', methods=['POST'])
def download_answers_docx():
    """
    Converts the provided list of questions and their generated answers into a .docx file
    and returns it for download.
    Formats bold text and bullet points from Markdown-like input.
    Also attempts to embed images found in the question text.
    """
    data = request.json
    answers_data = data.get('answers', []) # This should be a list of {'question', 'answer'}
    qb_filename = data.get('qbFilename', 'Generated_Answers')

    if not answers_data:
        return jsonify({'error': 'No answers provided to download.'}), 400

    document = Document()
    
    # Set document title
    document.add_heading(f'Generated Answers from {qb_filename}', level=0)
    document.add_paragraph("\n") # Add a blank line for spacing

    for item in answers_data:
        if 'question' in item and 'answer' in item:
            # Add question as a heading
            question_heading = document.add_heading('', level=2)
            
            # Process question text for images
            question_text_with_images = item['question']
            
            # Regex to find <img> tags and extract src attribute
            img_tag_pattern = re.compile(r'<img\s+src=["\']([^"\']+)["\'][^>]*>')
            
            last_idx = 0
            for match in img_tag_pattern.finditer(question_text_with_images):
                # Add text before the image
                text_before = question_text_with_images[last_idx:match.start()]
                if text_before:
                    question_heading.add_run(text_before)
                
                img_url_relative = match.group(1) # Get the relative URL
                
                # Convert relative URL to absolute path on the server
                # Assumes images are in static/temp_images
                img_path = os.path.join(app.root_path, img_url_relative.lstrip('/')) # Remove leading /
                
                if os.path.exists(img_path):
                    try:
                        # docx library requires specific image types, convert if necessary
                        img_extension = os.path.splitext(img_path)[1].lower().lstrip('.')
                        if img_extension not in ['png', 'jpeg', 'jpg', 'gif']:
                            # Try to convert to PNG
                            img = Image.open(img_path)
                            temp_png_path = img_path + ".png"
                            img.save(temp_png_path, format='PNG')
                            img_path = temp_png_path # Use the converted path
                            
                        # Add image to document, scale for better fitting
                        question_heading.add_run().add_picture(img_path, width=Inches(4)) # Adjust width as needed
                    except Exception as e:
                        print(f"Error embedding image {img_path} in DOCX: {e}")
                        question_heading.add_run(f"[Image: {img_url_relative} - failed to embed]")
                else:
                    question_heading.add_run(f"[Image: {img_url_relative} - file not found]")
                
                last_idx = match.end()
            
            # Add any remaining text after the last image
            remaining_text = question_text_with_images[last_idx:]
            if remaining_text:
                question_heading.add_run(remaining_text)

            # Process answer content for formatting
            answer_lines = item['answer'].split('\n')
            for line in answer_lines:
                stripped_line = line.strip()
                if not stripped_line:
                    continue # Skip empty lines

                # Handle Sources line explicitly
                if stripped_line.startswith('**Sources:**'):
                    p = document.add_paragraph()
                    run = p.add_run(stripped_line)
                    run.bold = True
                    p.paragraph_format.space_before = Inches(0.1) # Small space before sources
                    p.paragraph_format.space_after = Inches(0.05)
                elif stripped_line.startswith('- '): # Likely a bullet point or a source reference
                    p = document.add_paragraph(style='List Bullet')
                    run = p.add_run(stripped_line[2:])
                elif stripped_line.startswith('* '): # Another common bullet point style
                    p = document.add_paragraph(style='List Bullet')
                    run = p.add_run(stripped_line[2:])
                elif '**' in stripped_line: # Handle inline bolding
                    p = document.add_paragraph()
                    parts = re.split(r'(\*\*.*?\*\*)', stripped_line) # Split by bold segments
                    for part in parts:
                        if part.startswith('**') and part.endswith('**'):
                            p.add_run(part[2:-2]).bold = True
                        else:
                            p.add_run(part)
                else:
                    document.add_paragraph(stripped_line) # Regular paragraph

            document.add_page_break() # Add a page break after each Q&A

    # Save the document to a BytesIO object
    byte_stream = BytesIO()
    document.save(byte_stream)
    byte_stream.seek(0) # Go to the beginning of the stream

    # Send the file as a response
    return send_file(
        byte_stream,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        as_attachment=True,
        download_name=f'{os.path.splitext(qb_filename)[0]}_answers.docx' # Clean filename
    )


@app.route('/clear_all_data', methods=['POST'])
def clear_all_data_placeholder():
    """
    Placeholder for clearing data. As per user's request, no Firebase/persistent storage is used.
    This function simply resets the in-memory global variables and clears temp images.
    """
    global QUESTION_BANK_CONTENT, QUESTION_BANK_FILENAME, SLIDE_CONTENTS
    QUESTION_BANK_CONTENT = ""
    QUESTION_BANK_FILENAME = ""
    SLIDE_CONTENTS = {}
    print("In-memory data cleared (no persistent storage to clear).")

    # Also clear temporary images
    import shutil
    if os.path.exists(TEMP_IMAGE_DIR):
        shutil.rmtree(TEMP_IMAGE_DIR)
        os.makedirs(TEMP_IMAGE_DIR) # Recreate empty directory
        print(f"Cleared temporary image directory: {TEMP_IMAGE_DIR}")

    return jsonify({'message': 'All in-memory data and temporary images cleared successfully.'})


# Serve static files from the 'static' directory
@app.route('/static/<path:filename>')
def static_files(filename):
    # This route will serve files from the 'static' directory, including temp_images
    return send_file(os.path.join(app.root_path, 'static', filename))


# --- HTML Template ---
# This is a large multi-line string containing the entire HTML structure.
# It uses Jinja2-like syntax ({{ }}) for injecting Python variables.
HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Deep Learning Q&A Generator</title>
    <!-- Tailwind CSS CDN -->
    <script src="https://cdn.tailwindcss.com"></script>
    <style>
        body {
            font-family: 'Inter', sans-serif;
            background-color: #f4f7f6;
            color: #333;
            line-height: 1.6;
        }
        .container {
            max-width: 900px;
            margin: 2rem auto;
            padding: 1.5rem;
            background-color: #ffffff;
            border-radius: 12px;
            box-shadow: 0 10px 30px rgba(0, 0, 0, 0.08);
        }
        h1, h2 {
            color: #2c3e50;
            font-weight: 700;
        }
        .file-input-label {
            cursor: pointer;
            display: block;
            width: 100%;
            padding: 1rem;
            border: 2px dashed #a0aec0;
            border-radius: 8px;
            text-align: center;
            background-color: #edf2f7;
            transition: all 0.3s ease;
        }
        .file-input-label:hover {
            background-color: #e2e8f0;
            border-color: #4a5568;
        }
        .file-input-label input[type="file"] {
            display: none;
        }
        .loading-spinner {
            border-top-color: #3498db;
            -webkit-animation: spin 1s linear infinite;
            animation: spin 1s linear infinite;
        }
        @-webkit-keyframes spin {
            0% { -webkit-transform: rotate(0deg); }
            100% { -webkit-transform: rotate(360deg); }
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
        .progress-bar-container {
            width: 100%;
            background-color: #e0e0e0;
            border-radius: 5px;
            margin-top: 10px;
        }
        .progress-bar {
            width: 0%;
            height: 20px;
            background-color: #4CAF50;
            border-radius: 5px;
            text-align: center;
            color: white;
            line-height: 20px;
            transition: width 0.3s ease-in-out;
        }
        .message-modal {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background-color: rgba(0, 0, 0, 0.5);
            display: flex;
            justify-content: center;
            align-items: center;
            z-index: 1000;
        }
        .message-modal-content {
            background-color: white;
            padding: 30px;
            border-radius: 10px;
            box-shadow: 0 5px 15px rgba(0, 0, 0, 0.3);
            text-align: center;
            max-width: 400px;
            width: 90%;
        }
        .question-card {
            background-color: #fdfdfd;
            border: 1px solid #e2e8f0;
            border-radius: 8px;
            padding: 1rem;
            margin-bottom: 1rem;
            box-shadow: 0 2px 4px rgba(0, 0, 0, 0.05);
        }
        .question-card.answered {
            background-color: #e6fffa; /* Light teal for answered questions */
            border-color: #38b2ac; /* Teal border */
        }
        .answer-text {
            white-space: pre-wrap; /* Preserves formatting from AI response */
            word-wrap: break-word;
        }
        .pagination-btn {
            padding: 0.5rem 1rem;
            border-radius: 0.5rem;
            background-color: #e2e8f0;
            color: #2d3748;
            font-weight: 600;
            transition: background-color 0.2s;
        }
        .pagination-btn:hover:not(:disabled) {
            background-color: #cbd5e0;
        }
        .pagination-btn:disabled {
            opacity: 0.5;
            cursor: not-allowed;
        }
        .pagination-btn.active {
            background-color: #4299e1;
            color: white;
        }
        /* Style for extracted images within questions */
        .extracted-image {
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            margin: 10px 0;
            display: block; /* Ensures it takes full width and new line */
        }
    </style>
</head>
<body class="bg-gray-100 p-6"
      data-app-id='{{ app_id }}'
      data-firebase-config='{{ firebase_config_json }}'
      data-initial-auth-token='{{ initial_auth_token }}'>

    <div class="container">
        <h1 class="text-3xl font-bold text-center mb-6">Deep Learning Q&A Generator</h1>

        <!-- Question Bank Upload Section -->
        <div class="mb-8 p-6 border border-gray-200 rounded-lg shadow-sm bg-blue-50">
            <h2 class="text-xl font-semibold text-blue-800 mb-4">1. Upload Question Bank (TXT, PDF or PPTX)</h2>
            <label for="qbFile" class="file-input-label">
                <span id="qbFileName">Choose File(s)</span>
                <input type="file" id="qbFile" accept=".txt,.pdf,.pptx" multiple>
            </label>
            <button id="uploadQbBtn" class="mt-2 bg-blue-500 hover:bg-blue-600 text-white font-bold py-2 px-4 rounded-lg shadow transition duration-200 ease-in-out w-full flex items-center justify-center">
                <span id="qbUploadSpinner" class="hidden loading-spinner mr-2"></span>
                Upload Question Bank
            </button>
            <p id="qbStatus" class="mt-2 text-center text-sm text-gray-500"></p>
            <div class="progress-bar-container mt-2 hidden" id="qbProgressBarContainer">
                <div class="progress-bar" id="qbProgressBar">0%</div>
            </div>
            <p id="qbProcessingStatus" class="mt-2 text-center text-sm text-gray-700 hidden">Processing questions...</p>
            <div id="questionListSection" class="mt-6 hidden">
                <h3 class="text-lg font-semibold mb-3">Parsed Questions:</h3>
                <div id="questionList" class="max-h-60 overflow-y-auto border rounded p-3 bg-gray-50">
                    <!-- Questions will be loaded here -->
                </div>
                 <div class="flex justify-center items-center space-x-2 mt-4">
                    <button id="prevPageBtn" class="pagination-btn">Previous</button>
                    <span id="pageInfo" class="text-gray-700">Page 1 of 1</span>
                    <button id="nextPageBtn" class="pagination-btn">Next</button>
                </div>
            </div>
        </div>

        <!-- Lecture Slides Upload Section -->
        <div class="mb-8 p-6 border border-gray-200 rounded-lg shadow-sm bg-green-50">
            <h2 class="text-xl font-semibold text-green-800 mb-4">2. Upload Lecture Slides (PDF, PPTX, or TXT)</h2>
            <label for="slideFiles" class="file-input-label">
                <span id="slidesFileName">Choose File(s)</span>
                <input type="file" id="slideFiles" accept=".pdf,.pptx,.txt" multiple>
            </label>
            <button id="uploadSlidesBtn" class="mt-2 bg-green-500 hover:bg-green-600 text-white font-bold py-2 px-4 rounded-lg shadow transition duration-200 ease-in-out w-full flex items-center justify-center">
                <span id="slidesUploadSpinner" class="hidden loading-spinner mr-2"></span>
                Upload Lecture Slides
            </button>
            <p id="slidesStatus" class="mt-2 text-center text-sm text-gray-500"></p>
             <div class="progress-bar-container mt-2 hidden" id="slidesProgressBarContainer">
                <div class="progress-bar" id="slidesProgressBar">0%</div>
            </div>
            <p id="slidesProcessingStatus" class="mt-2 text-center text-sm text-gray-700 hidden">Processing slides...</p>
            <div id="uploadedSlidesList" class="mt-6 max-h-40 overflow-y-auto border rounded p-3 bg-gray-50 hidden">
                <h3 class="text-lg font-semibold mb-2">Uploaded Slides:</h3>
                <ul id="slidesListItems" class="list-disc list-inside">
                    <!-- Uploaded slides will be listed here -->
                </ul>
            </div>
        </div>

        <!-- Answer Generation Section -->
        <div class="mb-8 p-6 border border-gray-200 rounded-lg shadow-sm bg-yellow-50">
            <h2 class="text-xl font-semibold text-yellow-800 mb-4">3. Generate Answers</h2>
            <button id="generateAllAnswersBtn" class="bg-indigo-600 hover:bg-indigo-700 text-white font-bold py-3 px-6 rounded-lg shadow transition duration-200 ease-in-out w-full flex items-center justify-center disabled:opacity-50 disabled:cursor-not-allowed" disabled>
                Generate Answers for All Questions
                <span id="generateAllSpinner" class="hidden ml-3 w-5 h-5 border-2 border-white border-t-transparent rounded-full animate-spin"></span>
            </button>
            <p id="generateAllStatus" class="mt-4 text-center text-sm text-gray-500"></p>
            <div class="progress-bar-container mt-4 hidden" id="overallProgressBarContainer">
                <div class="progress-bar" id="overallProgressBar">0%</div>
            </div>
            <div id="answerDisplaySection" class="mt-6 hidden">
                <h3 class="text-lg font-semibold mb-3">Generated Answers:</h3>
                <div id="answersContainer" class="max-h-96 overflow-y-auto border rounded p-3 bg-gray-50">
                    <!-- Answers will be loaded here -->
                </div>
            </div>
            <button id="downloadDocxBtn" class="mt-4 bg-teal-600 hover:bg-teal-700 text-white font-bold py-2 px-4 rounded-lg shadow transition duration-200 ease-in-out w-full" disabled>
                Download as DOCX
            </button>
        </div>

        <!-- Clear All Data Section -->
        <div class="mb-8 p-6 border border-gray-200 rounded-lg shadow-sm bg-red-50">
            <h2 class="text-xl font-semibold text-red-800 mb-4">Clear All User Data</h2>
            <p class="text-sm text-gray-700 mb-4">This will delete all uploaded Question Banks, Lecture Slides, and generated answers from the server's memory. This action cannot be undone for the current session.</p>
            <button id="clearAllDataBtn" class="bg-red-600 hover:bg-red-700 text-white font-bold py-2 px-4 rounded-lg shadow transition duration-200 ease-in-out w-full flex items-center justify-center">
                Clear All My Data
                <span id="clearDataSpinner" class="hidden ml-2 w-4 h-4 border-2 border-white border-t-transparent rounded-full animate-spin"></span>
            </button>
            <p id="clearDataStatus" class="mt-4 text-center text-sm text-gray-500"></p>
        </div>
    </div>

    <!-- Link to external JavaScript module -->
    <script type="module" src="{{ url_for('static', filename='script.js') }}"></script>
</body>
</html>
"""

if __name__ == '__main__':
    # For development, run in debug mode. In production, use a WSGI server like Gunicorn or Waitress.
    app.run(debug=True)

